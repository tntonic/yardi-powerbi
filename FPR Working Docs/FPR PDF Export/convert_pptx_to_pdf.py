#!/usr/bin/env python3
"""
Script to convert PowerPoint presentations to PDFs with specific section extraction
"""

import os
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF
from datetime import datetime

def extract_slide_as_image(slide, width=1920, height=1080):
    """Convert a PowerPoint slide to an image"""
    # Create a white background image
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    
    # Try to get a default font
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 40)
        title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 60)
        small_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 30)
    except:
        font = ImageFont.load_default()
        title_font = font
        small_font = font
    
    y_position = 100
    
    # Extract and draw slide title
    if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
        title_text = slide.shapes.title.text
        if title_text:
            # Draw title with word wrapping
            words = title_text.split()
            lines = []
            current_line = []
            for word in words:
                current_line.append(word)
                test_line = ' '.join(current_line)
                bbox = draw.textbbox((0, 0), test_line, font=title_font)
                if bbox[2] > width - 200:
                    if len(current_line) > 1:
                        current_line.pop()
                        lines.append(' '.join(current_line))
                        current_line = [word]
                    else:
                        lines.append(test_line)
                        current_line = []
            if current_line:
                lines.append(' '.join(current_line))
            
            for line in lines:
                draw.text((100, y_position), line, fill='black', font=title_font)
                y_position += 70
            y_position += 30
    
    # Extract and draw other text content
    content_lines = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
            text = shape.text.strip()
            if text:
                content_lines.append(text)
    
    # Draw content with word wrapping
    for content in content_lines[:10]:  # Limit to first 10 content blocks
        words = content.split()
        lines = []
        current_line = []
        for word in words:
            current_line.append(word)
            test_line = ' '.join(current_line)
            bbox = draw.textbbox((0, 0), test_line, font=small_font)
            if bbox[2] > width - 200:
                if len(current_line) > 1:
                    current_line.pop()
                    lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    lines.append(test_line[:100] + '...')
                    break
            if len(lines) >= 3:  # Limit each content block to 3 lines
                lines.append('...')
                break
        if current_line and len(lines) < 3:
            lines.append(' '.join(current_line))
        
        for line in lines:
            if y_position < height - 100:
                draw.text((120, y_position), line, fill='#333333', font=small_font)
                y_position += 35
        y_position += 20
        
        if y_position > height - 150:
            draw.text((120, height - 100), "... (content continues)", fill='#666666', font=small_font)
            break
    
    return img

def create_pdf_from_slides(presentations, output_filename, slide_ranges=None):
    """Create a PDF from PowerPoint presentations
    
    Args:
        presentations: List of (filename, Presentation object) tuples
        output_filename: Output PDF filename
        slide_ranges: Optional dict mapping presentation filename to list of slide indices
    """
    pdf_document = fitz.open()
    
    for pptx_file, prs in presentations:
        print(f"  Processing: {pptx_file}")
        
        # Determine which slides to include
        if slide_ranges and pptx_file in slide_ranges:
            slides_to_process = slide_ranges[pptx_file]
        else:
            slides_to_process = range(len(prs.slides))
        
        for slide_idx in slides_to_process:
            if slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                
                # Convert slide to image
                img = extract_slide_as_image(slide)
                
                # Convert PIL Image to bytes
                img_bytes = io.BytesIO()
                img.save(img_bytes, format='PNG')
                img_bytes.seek(0)
                
                # Create a new page in the PDF with the image
                img_data = img_bytes.getvalue()
                img_doc = fitz.open(stream=img_data, filetype="png")
                
                # Get the first (and only) page from the image document
                pdfbytes = img_doc.convert_to_pdf()
                img_pdf = fitz.open("pdf", pdfbytes)
                pdf_document.insert_pdf(img_pdf)
                
                img_pdf.close()
                img_doc.close()
    
    # Save the PDF
    pdf_document.save(output_filename)
    pdf_document.close()
    print(f"  ✓ Created: {output_filename}")

def main():
    """Main execution function"""
    print("\n" + "="*60)
    print("PowerPoint to PDF Converter")
    print("="*60)
    
    # List of PowerPoint files
    pptx_files = [
        "FPR Fund 2 Q1 2025 - vF.pptx",
        "FPR Fund 2 Q2 2025 - vF.pptx",
        "FPR Fund 3 - Q1 2025.pptx",
        "FPR Fund 3 - Q2 2025.pptx"
    ]
    
    # Load all presentations
    presentations = []
    for pptx_file in pptx_files:
        if os.path.exists(pptx_file):
            print(f"Loading: {pptx_file}")
            prs = Presentation(pptx_file)
            presentations.append((pptx_file, prs))
        else:
            print(f"Warning: {pptx_file} not found")
    
    if not presentations:
        print("Error: No PowerPoint files found")
        return
    
    # Task 1: Combine all PowerPoints into one PDF
    print("\n>>> Task 1: Creating combined PDF with all presentations...")
    create_pdf_from_slides(presentations, "FPR_Combined_All_Funds_2025.pdf")
    
    # Task 2: Extract Fund Occupancy sections
    print("\n>>> Task 2: Creating Fund Occupancy PDF...")
    occupancy_ranges = {
        # Fund 2 Q1: Main occupancy section is slides 12-18
        "FPR Fund 2 Q1 2025 - vF.pptx": [11, 12, 13, 14, 15, 16, 17],  # 0-indexed
        # Fund 2 Q2: Main occupancy section is slides 11-18
        "FPR Fund 2 Q2 2025 - vF.pptx": [10, 11, 12, 13, 14, 15, 16, 17],
        # Fund 3 Q1: Main occupancy section is slides 12-16
        "FPR Fund 3 - Q1 2025.pptx": [11, 12, 13, 14, 15],
        # Fund 3 Q2: Main occupancy section is slides 16-21
        "FPR Fund 3 - Q2 2025.pptx": [15, 16, 17, 18, 19, 20]
    }
    create_pdf_from_slides(presentations, "FPR_Fund_Occupancy_Fund2_Fund3.pdf", occupancy_ranges)
    
    # Task 3: Extract Leasing Spreads & Downtime sections
    print("\n>>> Task 3: Creating Leasing Spreads & Downtime PDF...")
    leasing_ranges = {
        # Fund 2 Q1: Leasing section is slides 19-30
        "FPR Fund 2 Q1 2025 - vF.pptx": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29],
        # Fund 2 Q2: Leasing section is slides 19-29
        "FPR Fund 2 Q2 2025 - vF.pptx": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28],
        # Fund 3 Q1: Leasing section is slides 17-27
        "FPR Fund 3 - Q1 2025.pptx": [16, 17, 18, 19, 20, 21, 22, 23, 24, 25, 26],
        # Fund 3 Q2: Leasing section is slides 22-32
        "FPR Fund 3 - Q2 2025.pptx": [21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31]
    }
    create_pdf_from_slides(presentations, "FPR_Leasing_Spreads_Downtime.pdf", leasing_ranges)
    
    print("\n" + "="*60)
    print("PDF Generation Complete!")
    print("="*60)
    print("\nCreated files:")
    print("1. FPR_Combined_All_Funds_2025.pdf - All presentations combined")
    print("2. FPR_Fund_Occupancy_Fund2_Fund3.pdf - Fund Occupancy sections only")
    print("3. FPR_Leasing_Spreads_Downtime.pdf - Leasing Spreads & Downtime sections only")
    
    # Check file sizes
    print("\nFile sizes:")
    for filename in ["FPR_Combined_All_Funds_2025.pdf", 
                     "FPR_Fund_Occupancy_Fund2_Fund3.pdf", 
                     "FPR_Leasing_Spreads_Downtime.pdf"]:
        if os.path.exists(filename):
            size = os.path.getsize(filename)
            print(f"  {filename}: {size:,} bytes ({size/1024/1024:.1f} MB)")

if __name__ == "__main__":
    main()