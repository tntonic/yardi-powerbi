#!/usr/bin/env python3
"""
Script to examine PowerPoint presentations and identify slide sections
"""

from pptx import Presentation
import os

def examine_presentation(filename):
    """Examine a PowerPoint presentation and extract slide information"""
    print(f"\n{'='*60}")
    print(f"Examining: {filename}")
    print('='*60)
    
    prs = Presentation(filename)
    
    slides_info = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_title = ""
        slide_text = []
        
        # Try to get slide title
        if slide.shapes.title:
            slide_title = slide.shapes.title.text if hasattr(slide.shapes.title, 'text') else ""
        
        # Get all text from slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text and text != slide_title:
                    slide_text.append(text)
        
        # Store slide info
        slide_info = {
            'slide_number': idx,
            'title': slide_title,
            'has_title': bool(slide_title),
            'text_preview': ' | '.join(slide_text[:3]) if slide_text else ""
        }
        slides_info.append(slide_info)
        
        # Print slide info
        print(f"\nSlide {idx}:")
        print(f"  Title: {slide_title if slide_title else '(No title)'}")
        if slide_text:
            print(f"  Content preview: {slide_info['text_preview'][:100]}...")
    
    return slides_info

def find_sections(slides_info, keywords):
    """Find slides containing specific keywords"""
    matching_slides = []
    for slide in slides_info:
        title_lower = slide['title'].lower()
        text_lower = slide['text_preview'].lower()
        
        for keyword in keywords:
            if keyword.lower() in title_lower or keyword.lower() in text_lower:
                matching_slides.append(slide)
                break
    
    return matching_slides

# Main execution
if __name__ == "__main__":
    # List all PowerPoint files
    pptx_files = [
        "FPR Fund 2 Q1 2025 - vF.pptx",
        "FPR Fund 2 Q2 2025 - vF.pptx",
        "FPR Fund 3 - Q1 2025.pptx",
        "FPR Fund 3 - Q2 2025.pptx"
    ]
    
    all_slides = {}
    
    # Examine each presentation
    for pptx_file in pptx_files:
        if os.path.exists(pptx_file):
            slides_info = examine_presentation(pptx_file)
            all_slides[pptx_file] = slides_info
    
    # Search for specific sections
    print("\n" + "="*60)
    print("SEARCHING FOR SPECIFIC SECTIONS")
    print("="*60)
    
    # Search for Fund Occupancy
    print("\n>>> Searching for 'Fund Occupancy' sections:")
    for pptx_file, slides in all_slides.items():
        occupancy_slides = find_sections(slides, ['occupancy', 'fund occupancy'])
        if occupancy_slides:
            print(f"\n{pptx_file}:")
            for slide in occupancy_slides:
                print(f"  - Slide {slide['slide_number']}: {slide['title']}")
    
    # Search for Leasing Spreads & Downtime
    print("\n>>> Searching for 'Leasing Spreads & Downtime' sections:")
    for pptx_file, slides in all_slides.items():
        leasing_slides = find_sections(slides, ['leasing spread', 'downtime', 'leasing spreads'])
        if leasing_slides:
            print(f"\n{pptx_file}:")
            for slide in leasing_slides:
                print(f"  - Slide {slide['slide_number']}: {slide['title']}")
    
    # Summary
    print("\n" + "="*60)
    print("SUMMARY")
    print("="*60)
    for pptx_file, slides in all_slides.items():
        print(f"{pptx_file}: {len(slides)} slides")